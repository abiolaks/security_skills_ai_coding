#!/usr/bin/env python3
"""
Build "Secure by Default — Security in AI Agentic Workflows" as a PowerPoint
deck that mirrors the HTML design (slides/codeguard-agentic-workflow.html):
Wragby brand (Figtree, #EE1F2F red, #00ADEF cyan), light theme, SVG icons,
dark code panels with the repo's Python examples.

Requires the icon PNGs rendered to /tmp/icons_png (see inline note) and the
Figtree + JetBrains Mono fonts installed. Logo at slides/wragby-logo.png.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Brand tokens (mirror the HTML)
# ----------------------------------------------------------------------------
INK      = RGBColor(0x0E, 0x0F, 0x11)
INK2     = RGBColor(0x3C, 0x40, 0x47)
MUTED    = RGBColor(0x6A, 0x70, 0x78)
FAINT    = RGBColor(0x9A, 0xA1, 0xAB)
LINE     = RGBColor(0xE7, 0xE9, 0xEC)
SURFACE  = RGBColor(0xF6, 0xF7, 0xF8)
SURFACE2 = RGBColor(0xEE, 0xF0, 0xF2)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RED      = RGBColor(0xEE, 0x1F, 0x2F)
RED_SOFT = RGBColor(0xFD, 0xEC, 0xEE)
CYAN     = RGBColor(0x00, 0xAD, 0xEF)
CYAN_SOFT= RGBColor(0xE7, 0xF6, 0xFC)
GREEN    = RGBColor(0x15, 0x9A, 0x4B)
GREEN_SOFT=RGBColor(0xE9, 0xF7, 0xEF)
AMBER    = RGBColor(0xC0, 0x7D, 0x00)
AMBER_SOFT=RGBColor(0xFB, 0xF3, 0xE3)
CODE_BG  = RGBColor(0x0F, 0x11, 0x14)

# code token colors
C_DEF  = RGBColor(0xDB, 0xE1, 0xE8)
C_MUT  = RGBColor(0x8B, 0x94, 0x9E)
C_KEY  = RGBColor(0xFF, 0x7B, 0x72)
C_STR  = RGBColor(0x7E, 0xE7, 0x87)
C_FN   = RGBColor(0x79, 0xC0, 0xFF)
C_NUM  = RGBColor(0xF2, 0xCC, 0x60)
C_DANG = RGBColor(0xFF, 0xB4, 0xBB)
C_SAFE = RGBColor(0x8C, 0xE2, 0xAB)

FONT = "Figtree"
MONO = "JetBrains Mono"

SLIDE_W = 13.333
SLIDE_H = 7.5
MX = 0.6                      # page margin

ICON_DIR = "/tmp/icons_png"
LOGO = os.path.join(os.path.dirname(__file__), "..", "slides", "wragby-logo.png")

COLOR_HEX = {
    "red":   "#ee1f2f", "cyan": "#00adef", "green": "#159a4b",
    "amber": "#c07d00", "gray": "#9aa1ab", "ink": "#3c4047",
}
COLOR_RGB = {"red": RED, "cyan": CYAN, "green": GREEN, "amber": AMBER, "gray": FAINT, "ink": INK2}
SOFT_RGB  = {"red": RED_SOFT, "cyan": CYAN_SOFT, "green": GREEN_SOFT, "amber": AMBER_SOFT}

TOKEN = {"d": C_DEF, "c": C_MUT, "k": C_KEY, "s": C_STR, "f": C_FN, "n": C_NUM}

# ----------------------------------------------------------------------------
# low-level helpers
# ----------------------------------------------------------------------------
def _noshadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def solid(shape, rgb, line_rgb=None, line_w=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_w if line_w else 1)
    _noshadow(shape)
    return shape


def rect(slide, l, t, w, h, rgb, radius=None, border=None, border_w=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    if radius:
        try:
            sp.adjustments[0] = max(0.0, min(0.5, radius / min(w, h)))
        except Exception:
            pass
    solid(sp, rgb, border, border_w)
    return sp


def set_spacing(run, val_100pt):
    run._r.get_or_add_rPr().set("spc", str(val_100pt))


def txt(slide, l, t, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True, align=None):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", align if align is not None else PP_ALIGN.LEFT)
        if p.get("sp") is not None:
            para.space_after = Pt(p["sp"])
        if p.get("line_spacing"):
            para.line_spacing = p["line_spacing"]
        runs = p.get("runs") or [p]
        for r in runs:
            run = para.add_run()
            run.text = r.get("t", "")
            f = run.font
            f.name = r.get("font", p.get("font", FONT))
            f.size = Pt(r.get("s", p.get("s", 12)))
            f.bold = r.get("b", p.get("b", False))
            f.color.rgb = r.get("c", p.get("c", INK))
            if r.get("spc"):
                set_spacing(run, r["spc"])
    return box


def _p(text, **kw):
    d = dict(kw)
    d["t"] = text
    return d


def icon_png(name, color):
    p = os.path.join(ICON_DIR, f"{name}__{color}.png")
    assert os.path.exists(p), f"missing icon {name}__{color}"
    return p


def add_icon(slide, l, t, size, name, color):
    slide.shapes.add_picture(icon_png(name, color), Inches(l), Inches(t), Inches(size), Inches(size))


def icon_tile(slide, l, t, size, name, color, tile_rgb=None):
    rect(slide, l, t, size, size, tile_rgb if tile_rgb is not None else SOFT_RGB[color], radius=size * 0.26)
    add_icon(slide, l + size * 0.22, t + size * 0.22, size * 0.56, name, color)


# ----------------------------------------------------------------------------
# code panel
# ----------------------------------------------------------------------------
LINE_H = 0.205

def add_code(slide, l, t, w, lines):
    """lines: list of dicts {runs:[(text, token)], kind:'danger'|'safe'|None}
    Returns bottom Y of the panel."""
    pad_t = 0.17
    pad_x = 0.22
    h = pad_t * 2 + len(lines) * LINE_H
    rect(slide, l, t, w, h, CODE_BG, radius=0.09)
    y = t + pad_t
    for ln in lines:
        kind = ln.get("kind")
        if kind:
            tint = RED if kind == "danger" else GREEN
            alpha_rgb = None
            bar = rect(slide, l + 0.02, y - 0.025, w - 0.04, LINE_H, tint, radius=0.02)
            _set_alpha(bar, 16)
        runs = []
        for text, tok in ln["runs"]:
            color = C_DANG if kind == "danger" else (C_SAFE if kind == "safe" else TOKEN.get(tok, C_DEF))
            runs.append(_p(text, c=color, font=MONO, s=10.5))
        txt(slide, l + pad_x, y, w - pad_x * 2, LINE_H, [{"runs": runs}])
        y += LINE_H
    return t + h


def _set_alpha(shape, pct):
    sp = shape._element
    sf = sp.spPr.find(qn("a:solidFill"))
    clr = sf.find(qn("a:srgbClr"))
    for a in clr.findall(qn("a:alpha")):
        clr.remove(a)
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))


# ----------------------------------------------------------------------------
# chrome
# ----------------------------------------------------------------------------
FOOTER = "Secure by Default — Security in AI Agentic Workflows"

def chrome(slide, section_tag, page):
    slide.shapes.add_picture(LOGO, Inches(MX), Inches(0.30), height=Inches(0.27))
    # right-aligned section tag
    txt(slide, 8.7, 0.32, 13.333 - MX - 8.7, 0.3,
        [_p(section_tag, s=10.5, b=True, c=RED, spc=160, align=PP_ALIGN.RIGHT)])
    # footer
    rect(slide, 0, 6.98, SLIDE_W, 0.012, LINE)
    txt(slide, MX, 7.10, 8.0, 0.25, [_p(FOOTER, s=9, c=MUTED)])
    txt(slide, 11.3, 7.10, 1.43, 0.25, [_p(f"{page:02d} / 14", s=9, c=FAINT, align=PP_ALIGN.RIGHT)])


def header(slide, eyebrow, title, lede):
    txt(slide, MX, 1.00, 11.0, 0.3, [_p(eyebrow.upper(), s=11, b=True, c=RED, spc=170)])
    txt(slide, MX, 1.28, 11.9, 0.75, [_p(title, s=30, b=True, c=INK)])
    if lede:
        txt(slide, MX, 2.06, 11.5, 0.55, [_p(lede, s=13.5, c=MUTED, line_spacing=1.15)])


# ----------------------------------------------------------------------------
# card
# ----------------------------------------------------------------------------
def card(slide, l, t, w, h, icon=None, icon_color="red", title=None, lines=None,
         accent=None, title_s=13.5):
    rect(slide, l, t, w, h, SURFACE, radius=0.09, border=LINE, border_w=1)
    if accent:
        rect(slide, l, t, 0.045, h, accent)
    yy = t + 0.18
    if icon:
        icon_tile(slide, l + 0.18, t + 0.16, 0.40, icon, icon_color)
        yy = t + 0.18 + 0.40 + 0.14
    if title:
        txt(slide, l + 0.18, yy, w - 0.36, 0.3, [_p(title, s=title_s, b=True, c=INK)])
        yy += 0.30
    if lines:
        paras = []
        for ln in lines:
            if isinstance(ln, dict):
                paras.append(ln)
            else:
                paras.append(_p(ln, s=11, c=INK2, sp=5))
        txt(slide, l + 0.18, yy, w - 0.36, h - (yy - t) - 0.12, paras)


def bullets(marker_color, items, s=11):
    paras = []
    for it in items:
        paras.append({
            "runs": [
                _p("▪  ", c=marker_color, s=s, b=True),
                _p(it, c=INK2, s=s),
            ],
            "sp": 6,
            "line_spacing": 1.05,
        })
    return paras


# ============================================================================
# BUILD
# ============================================================================
def build(out_path):
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]

    # ------------------------------------------------------------ 1 · TITLE
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(LOGO, Inches(MX + 0.12), Inches(0.62), height=Inches(0.34))
    txt(s, MX + 0.12, 1.75, 11.0, 0.3,
        [_p("WRAGBY BUSINESS SOLUTIONS & TECHNOLOGIES · AI & SECURITY ENGINEERING", s=11, b=True, c=RED, spc=150)])
    txt(s, MX + 0.10, 2.10, 11.5, 1.3, [_p("Secure by Default", s=54, b=True, c=INK)])
    txt(s, MX + 0.12, 3.35, 11.5, 0.8, [_p("Security in ", s=34, b=True, c=INK),
                                         _p("AI Agentic Workflows", s=34, b=True, c=RED)])
    txt(s, MX + 0.12, 4.35, 9.5, 0.9,
        [_p("How ", s=16, c=MUTED), _p("Project CodeGuard", s=16, b=True, c=INK2),
         _p(" and the agentic workflow stop vulnerabilities before they reach production — "
            "and why \u201cprompt and pray\u201d is a security risk.", s=16, c=MUTED)],
        )
    txt(s, MX + 0.12, 5.55, 9.0, 0.3,
        [_p("Presenter · ", s=13, c=MUTED), _p("Your Name", s=13, b=True, c=INK)])
    txt(s, MX + 0.12, 5.85, 9.0, 0.3,
        [_p("github.com/abiolaks/security_skills_ai_coding", s=12.5, c=INK2, font=MONO)])

    # ----------------------------------------------------- 2 · SPEED TRAP
    s = prs.slides.add_slide(blank)
    chrome(s, "The Problem · 01", 2)
    header(s, "The Problem", "The AI Coding Speed Trap",
           "AI agents ship code 2–5× faster — and that speed amplifies every bad habit. "
           "The code below looks right, runs, and still ships three bugs.")
    cards = [
        ("i-box", "red", "Hallucinated SDK methods",
         [_p("openai.Embedding.generate()", c=INK2, font=MONO, s=10.5),
          _p(" does not exist — it is ", c=INK2, s=11),
          _p("client.embeddings.create()", c=INK2, font=MONO, s=10.5),
          _p(".", c=INK2, s=11)]),
        ("i-lock", "red", "Hardcoded API keys",
         [_p("\u201csk-proj-…\u201d", c=INK2, font=MONO, s=10.5),
          _p(" lands in git history. One push and the key is revoked.", c=INK2, s=11)]),
        ("i-bug", "red", "Skipped edge cases",
         [_p("No rate limits, no input validation, no error handling. A 1 MB payload burns your API budget.", c=INK2, s=11)]),
    ]
    for i, (ic, cc, ti, li) in enumerate(cards):
        card(s, MX + i * 4.18, 2.62, 3.9, 1.62, icon=ic, icon_color=cc, title=ti, lines=li)

    code_lines = [
        {"runs": [("# 12 lines, 5 problems", "c")]},
        {"runs": [("from ", "k"), ("fastapi", "d"), (" import ", "k"), ("FastAPI, Request", "d")]},
        {"runs": [("from ", "k"), ("openai", "d"), (" import ", "k"), ("Embedding", "d"),
                  ("    # ← hallucinated import", "c")], "kind": "danger"},
        {"runs": [("from ", "k"), ("pinecone", "d"), (" import ", "k"), ("Pinecone", "d")]},
        {"runs": [("", "d")]},
        {"runs": [("app = FastAPI()", "d")]},
        {"runs": [("API_KEY = ", "d"), ("\"sk-proj-abc123\"", "s"),
                  ("    # ← hardcoded secret", "c")], "kind": "danger"},
        {"runs": [("pc = Pinecone(api_key=API_KEY)", "d")]},
        {"runs": [("", "d")]},
        {"runs": [("@app.post", "f"), ("(", "d"), ("\"/rag/query\"", "s"), (")", "d")]},
        {"runs": [("async def ", "k"), ("rag_query", "f"), ("(request: Request):", "d")]},
        {"runs": [("    query = body[", "d"), ("\"q\"", "s"), ("]", "d"),
                  ("       # ← no validation", "c")], "kind": "danger"},
        {"runs": [("    embedding = Embedding.generate(", "d"),
                  ("   # ← hallucinated method", "c")], "kind": "danger"},
        {"runs": [("        model=", "d"), ("\"text-embedding-3-small\"", "s"),
                  (", input=query)", "d")]},
        {"runs": [("    results = index.query(vector=embedding, top_k=", "d"), ("5", "n"), (")", "d")]},
        {"runs": [("    return ", "k"), ("{\"results\": results}", "d")]},
    ]
    add_code(s, MX, 4.42, 12.13, code_lines)

    # ----------------------------------------------------- 3 · SCALE
    s = prs.slides.add_slide(blank)
    chrome(s, "Evidence · 02", 3)
    header(s, "The Evidence", "The Scale of the Problem",
           "Cisco's controlled study: 2,717 prompts × 2 agents (baseline vs CodeGuard) × 9 languages, GPT-5.")
    stats = [("415", "security findings\nwithout CodeGuard", "red"),
             ("264", "security findings\nwith CodeGuard", "cyan"),
             ("36.4%", "reduction in vulnerabilities\np < 0.05 across benchmarks", "green")]
    for i, (val, lab, cc) in enumerate(stats):
        l = MX + i * 4.18
        rect(s, l, 2.62, 3.9, 1.7, SURFACE, radius=0.09, border=LINE)
        txt(s, l + 0.22, 2.82, 3.5, 0.7, [_p(val, s=40, b=True, c=COLOR_RGB[cc])])
        txt(s, l + 0.22, 3.60, 3.5, 0.65, [_p(lab, s=12, c=MUTED, line_spacing=1.1)])
    card(s, MX, 4.55, 5.95, 1.55, icon="i-target", icon_color="cyan",
         title="Hardest benchmark — SecurityEval",
         lines=[_p("59.1% fewer findings on the hardest security benchmark.", s=11.5, c=INK2)])
    card(s, MX + 6.18, 4.55, 5.95, 1.55, icon="i-trend", icon_color="green",
         title="Clean snippets improve",
         lines=[_p("68.6% → 85.1% of generated snippets with no security findings.", s=11.5, c=INK2)])
    txt(s, MX, 6.30, 11.5, 0.3,
        [_p("Source: Cisco / CoSAI — \u201cCan Security Rules Make AI Generated Code Safer?\u201d", s=10, c=FAINT)])

    # ----------------------------------------------------- 4 · OLD WAY
    s = prs.slides.add_slide(blank)
    chrome(s, "The Problem · 03", 4)
    header(s, "The Problem", "The Old Way: Prompt & Pray",
           "Most teams treat AI coding agents like a chatbot — unstructured, unverified, insecure.")
    flow = [("i-terminal", "Prompt", None), ("i-sparkles", "AI writes code", None),
            ("i-x-circle", "Ship?", "red")]
    fx = MX
    for i, (ic, t, cc) in enumerate(flow):
        rect(s, fx, 2.60, 3.7, 1.15, RED_SOFT if cc else SURFACE, radius=0.09,
             border=RED if cc else LINE)
        add_icon(s, fx + 1.62, 2.72, 0.30, ic, cc or "red")
        txt(s, fx, 3.12, 3.7, 0.3, [_p(t, s=13, b=True, c=INK, align=PP_ALIGN.CENTER)])
        if i < 2:
            add_icon(s, fx + 3.70, 2.88, 0.30, "i-arrow", "gray")
        fx += 4.18
    card(s, MX, 4.10, 5.95, 2.45, icon="i-x-circle", icon_color="red", title="What's missing",
         lines=bullets(RED, ["No specification — building the wrong thing",
                             "No security review — vulnerabilities ship silently",
                             "No standards check — inconsistent code",
                             "No testing discipline — regressions go unnoticed"], s=11.5))
    card(s, MX + 6.18, 4.10, 5.95, 2.45, icon="i-alert", icon_color="red", title="The result",
         lines=bullets(RED, ["SQL injection in production",
                             "Hardcoded API keys in git history",
                             "Weak hashing on passwords",
                             "Missing access controls",
                             "Nobody knows what the code should do"], s=11.5))
    txt(s, MX, 6.55, 12.13, 0.3,
        [_p("Speed without structure = ", s=14, b=True, c=INK),
         _p("amplified insecurity", s=14, b=True, c=RED)], align=PP_ALIGN.CENTER)

    # ----------------------------------------------------- 5 · AGENTIC WORKFLOW
    s = prs.slides.add_slide(blank)
    chrome(s, "The Solution · 04", 5)
    header(s, "The Solution", "The Agentic Workflow",
           "A structured pipeline that treats AI as an engineering teammate — not a magic oracle.")
    flow = [("i-flame", "Grill", "Stress-test the idea"), ("i-file", "To-spec", "Testable spec"),
            ("i-clipboard", "To-tickets", "Small work items"),
            ("i-terminal", "Implement", "TDD + security"), ("i-search", "Review", "Standards + spec")]
    fx = MX
    for i, (ic, t, d) in enumerate(flow):
        rect(s, fx, 2.58, 2.2, 1.5, SURFACE, radius=0.09, border=LINE)
        add_icon(s, fx + 0.93, 2.70, 0.30, ic, "red")
        txt(s, fx, 3.10, 2.2, 0.28, [_p(t, s=12.5, b=True, c=INK, align=PP_ALIGN.CENTER)])
        txt(s, fx + 0.08, 3.40, 2.04, 0.5, [_p(d, s=9.5, c=MUTED, align=PP_ALIGN.CENTER)])
        if i < 4:
            add_icon(s, fx + 2.20, 3.05, 0.26, "i-arrow", "gray")
        fx += 2.44
    labels = [("Structure", "Every phase has a purpose. Every output is verified before the next begins."),
              ("Quality", "TDD red→green bakes edge cases in. Typechecking runs throughout."),
              ("Correctness", "Two-axis review: standards (conventions) and spec (did it build the right thing).")]
    for i, (ti, de) in enumerate(labels):
        card(s, MX + i * 4.18, 4.35, 3.9, 1.45, title=ti, lines=[_p(de, s=11, c=INK2)])
    txt(s, MX, 6.05, 12.13, 0.3,
        [_p("This solves structure, quality and correctness — ", s=12.5, c=MUTED),
         _p("but there is still a gap.", s=12.5, b=True, c=INK)])

    # ----------------------------------------------------- 6 · SECURITY GAP
    s = prs.slides.add_slide(blank)
    chrome(s, "The Gap · 05", 6)
    header(s, "The Gap", "The Security Gap in Code Review",
           "Code review checks standards and spec. It does not check for SQL injection, hardcoded secrets, or weak crypto.")
    card(s, MX, 2.62, 5.95, 2.9, icon="i-check-circle", icon_color="green", title="What code review catches",
         lines=bullets(GREEN, ["Mysterious names, duplication, feature envy",
                               "Missing spec requirements, scope creep",
                               "Broken conventions, refactoring smells"], s=11.5))
    card(s, MX + 6.18, 2.62, 5.95, 2.9, icon="i-x-circle", icon_color="red", title="What code review misses",
         lines=bullets(RED, ["cursor.execute(f\"SELECT * FROM {table}\")",
                             "API_KEY = \"sk-live-…\"",
                             "hashlib.md5(password).hexdigest()",
                             "subprocess.run(cmd, shell=True)",
                             "Missing CSRF tokens, no CSP, root containers"], s=11))
    txt(s, MX, 5.75, 12.13, 0.4, [_p("Standards + Spec ", s=20, b=True, c=INK),
                                   _p("≠", s=20, b=True, c=RED),
                                   _p(" Security", s=20, b=True, c=INK)], align=PP_ALIGN.CENTER)
    txt(s, MX, 6.25, 12.13, 0.3, [_p("Security is a third axis that deserves its own dedicated review.", s=12.5, c=MUTED)],
        align=PP_ALIGN.CENTER)

    # ----------------------------------------------------- 7 · CODEGUARD
    s = prs.slides.add_slide(blank)
    chrome(s, "The Solution · 06", 7)
    header(s, "The Solution", "Enter Project CodeGuard",
           "An open-source framework by CoSAI / Cisco — 23 OWASP-backed security rules embedded into AI coding workflows.")
    rules = ["Cryptography", "Input Validation", "Authentication", "Authorization", "API Security",
             "Session Mgmt", "XSS / CSRF", "Data Storage", "Logging", "Supply Chain",
             "File Handling", "Docker / K8s", "IaC Security", "SSRF Prevention", "MFA / OAuth",
             "Hardcoded Secrets"]
    # rule tags as a wrapped paragraph of pill-like runs (simple text)
    tag_paras = []
    line = []
    txt(s, MX, 2.62, 12.13, 1.0, [_p("  ·  ".join(rules), s=11.5, c=INK2, line_spacing=1.25)])
    card(s, MX, 3.62, 5.95, 2.2, icon="i-layers", icon_color="cyan", title="Two-phase protection",
         lines=bullets(CYAN, ["During generation — rules loaded into the agent's context, preventing vulnerabilities before they are introduced.",
                              "After generation — a full 23-rule audit of the diff catches what slipped through."], s=11))
    card(s, MX + 6.18, 3.62, 5.95, 2.2, icon="i-gauge", icon_color="green", title="Proven impact (2,717 prompts)",
         lines=bullets(GREEN, ["36.4% fewer vulnerabilities overall",
                               "59.1% on SecurityEval · 28.9% on CyberSecEval",
                               "Clean snippets 68.6% → 85.1% · p < 0.05"], s=11))
    txt(s, MX, 6.05, 12.13, 0.3,
        [_p("Model-agnostic — GPT-5, Claude, Gemini — across Cursor, Copilot, Codex, Claude Code, Windsurf and Pi.", s=10.5, c=FAINT)])

    # ----------------------------------------------------- 8 · PHASE 1
    s = prs.slides.add_slide(blank)
    chrome(s, "How it works · 07", 8)
    header(s, "How it works", "Phase 1: During Generation",
           "Security rules are loaded into the agent's context before it writes a single line.")
    card(s, MX, 2.60, 5.95, 2.5, icon="i-shield", icon_color="red", title="Always-on rules (every session)",
         lines=bullets(RED, ["Never hardcode credentials — env vars or secrets manager",
                             "Never use weak crypto — Argon2id, AES-256-GCM, TLS 1.3",
                             "Never concatenate into queries — parameterized SQL, no shell=True",
                             "Validate all untrusted input at trust boundaries"], s=10.5))
    card(s, MX + 6.18, 2.60, 5.95, 2.5, icon="i-target", icon_color="cyan", title="Context-scoped rules",
         lines=bullets(CYAN, ["Auth → MFA, OAuth PKCE, Argon2id, rate limiting",
                              "API → SSRF prevention, HTTPS, schema validation",
                              "Frontend → XSS sinks, CSP, CSRF tokens",
                              "Docker / K8s → non-root, pinned digests"], s=10.5))
    code_lines = [
        {"runs": [("import os", "k")]},
        {"runs": [("from ", "k"), ("openai", "d"), (" import ", "k"), ("OpenAI", "d"),
                  ("           # ← correct import", "c")]},
        {"runs": [("", "d")]},
        {"runs": [("client = OpenAI(api_key=os.environ[\"OPENAI_API_KEY\"])", "d")], "kind": "safe"},
        {"runs": [("", "d")]},
        {"runs": [("@app.post", "f"), ("(\"/rag/query\")", "d")]},
        {"runs": [("@limiter.limit", "f"), ("(\"30/minute\")", "d"), ("      # ← rate limited", "c")]},
        {"runs": [("async def ", "k"), ("rag_query", "f"), ("(request: Request):", "d")]},
        {"runs": [("    body = ", "k"), ("await", "k"), (" request.json()", "d")], "kind": "safe"},
        {"runs": [("    query = (body ", "k"), ("or", "k"), (" {}).get(\"q\", \"\").strip()   # ← validated", "d")], "kind": "safe"},
        {"runs": [("    if not", "k"), (" query ", "k"), ("or", "k"), (" len(query) > ", "d"), ("1000", "n"),
                  (":      # ← bounded", "c")], "kind": "safe"},
        {"runs": [("        raise", "k"), (" HTTPException(", "d"), ("400", "n"), (", \"Invalid query\")", "s")], "kind": "safe"},
        {"runs": [("", "d")]},
        {"runs": [("    embedding = client.embeddings.create(   # ← real method", "d")], "kind": "safe"},
        {"runs": [("        model=", "d"), ("\"text-embedding-3-small\"", "s"), (", input=query)", "d")]},
        {"runs": [("    return ", "k"), ("{\"results\": index.query(...)}", "d")]},
    ]
    add_code(s, MX, 5.28, 12.13, code_lines)

    # ----------------------------------------------------- 9 · PHASE 2
    s = prs.slides.add_slide(blank)
    chrome(s, "How it works · 08", 9)
    header(s, "How it works", "Phase 2: After Generation",
           "A dedicated security sub-agent audits every changed line against all 23 CodeGuard rules.")
    flow = [("i-git", "git diff", "main…HEAD"), ("i-list", "Match rules", "to changed files"),
            ("i-shield-check", "Security sub-agent", "23-rule audit"), ("i-gauge", "Report", "by severity")]
    fx = MX
    for i, (ic, t, d) in enumerate(flow):
        rect(s, fx, 2.60, 2.85, 1.5, SURFACE, radius=0.09, border=LINE)
        add_icon(s, fx + 1.25, 2.70, 0.30, ic, "red")
        txt(s, fx, 3.10, 2.85, 0.28, [_p(t, s=12, b=True, c=INK, align=PP_ALIGN.CENTER)])
        txt(s, fx + 0.08, 3.42, 2.69, 0.5, [_p(d, s=9.5, c=MUTED, align=PP_ALIGN.CENTER)])
        if i < 3:
            add_icon(s, fx + 2.85, 3.05, 0.26, "i-arrow", "gray")
        fx += 3.10
    sev = [("CRITICAL", RED, RED_SOFT, "Credential leak, RCE, SQL injection", "Hardcoded AWS keys, raw SQL concat, pickle.loads on user data"),
           ("HIGH", AMBER, AMBER_SOFT, "Auth bypass, data exposure", "Missing CSRF, IDOR-vulnerable queries, MD5 for passwords"),
           ("MEDIUM", INK2, SURFACE2, "Weak config, hardening", "Missing security headers, no CSP, running as root in Docker"),
           ("LOW", FAINT, SURFACE2, "Improvement", "Missing structured logging, no SBOM, no healthcheck")]
    for i, (sevname, sevc, sevsoft, sevt, sevb) in enumerate(sev):
        col, row = i % 2, i // 2
        l = MX + col * 6.18
        t = 4.35 + row * 1.28
        rect(s, l, t, 5.95, 1.1, SURFACE, radius=0.09, border=LINE)
        rect(s, l + 0.16, t + 0.16, 0.95, 0.28, sevsoft, radius=0.14)
        txt(s, l + 0.16, t + 0.195, 0.95, 0.2, [_p(sevname, s=9.5, b=True, c=sevc, align=PP_ALIGN.CENTER)])
        txt(s, l + 1.24, t + 0.15, 4.55, 0.3, [_p(sevt, s=11.5, b=True, c=INK)])
        txt(s, l + 1.24, t + 0.45, 4.55, 0.6, [_p(sevb, s=10, c=INK2, line_spacing=1.05)])
    txt(s, MX, 6.20, 12.13, 0.3,
        [_p("Runs alongside (not instead of) code review — a separate axis with dedicated focus.", s=10.5, c=FAINT)],
        align=PP_ALIGN.CENTER)

    # ----------------------------------------------------- 10 · COMPLETE WORKFLOW
    s = prs.slides.add_slide(blank)
    chrome(s, "The Workflow · 09", 10)
    header(s, "The Workflow", "The Complete Workflow",
           "Security baked into every phase — not bolted on at the end.")
    phases = [("i-flame", "amber", "Before code", "Grill stress-tests the idea. Spec defines secure requirements. Tickets scope work into safe increments.", "Security starts at design — auth, data classification, threat model."),
              ("i-shield", "cyan", "During code", "CodeGuard rules loaded into context. Always-on rules prevent secrets, weak crypto and injection.", "Agent writes secure code by default — not \u201cfix it in review.\u201d"),
              ("i-shield-check", "green", "After code", "Dual review: codeguard-review audits 23 rules; code-review checks standards and spec.", "Defense in depth — what one misses, the other catches.")]
    for i, (ic, cc, ti, de, note) in enumerate(phases):
        l = MX + i * 4.18
        rect(s, l, 2.62, 3.9, 3.1, SURFACE, radius=0.09, border=LINE)
        rect(s, l, 2.62, 3.9, 0.06, COLOR_RGB[cc])
        icon_tile(s, l + 0.2, 2.82, 0.40, ic, cc)
        txt(s, l + 0.2, 3.40, 3.5, 0.3, [_p(ti, s=14, b=True, c=INK)])
        txt(s, l + 0.2, 3.72, 3.5, 1.0, [_p(de, s=10.5, c=INK2, line_spacing=1.12)])
        txt(s, l + 0.2, 4.95, 3.5, 0.7, [_p(note, s=9.5, c=MUTED, line_spacing=1.1)])
    txt(s, MX, 6.05, 12.13, 0.4,
        [_p("Secure by default. ", s=18, b=True, c=GREEN),
         _p("Not \u201csecure if we remember.\u201d", s=18, c=MUTED)], align=PP_ALIGN.CENTER)

    # ----------------------------------------------------- 11 · BEFORE / AFTER
    s = prs.slides.add_slide(blank)
    chrome(s, "Proof · 10", 11)
    header(s, "Proof", "Real Example: Before & After",
           "SecurityEval benchmark — same prompt, same model (GPT-5). The only difference is CodeGuard.")
    # before panel
    rect(s, MX, 2.58, 5.95, 0.46, RED_SOFT, radius=0.09)
    add_icon(s, MX + 0.18, 2.70, 0.20, "i-x-circle", "red")
    txt(s, MX + 0.46, 2.70, 5.0, 0.3, [_p("Without CodeGuard", s=12, b=True, c=RED)])
    before = [
        {"runs": [("raw_data = request.get_data(as_text=", "d"), ("True", "k"), (")", "d")]},
        {"runs": [("decoded = base64.b64decode(raw_data)", "d")]},
        {"runs": [("deserialized = pickle.loads(decoded)", "d")], "kind": "danger"},
        {"runs": [("# ← arbitrary code execution (CWE-502)", "c")], "kind": "danger"},
        {"runs": [("return ", "k"), ("deserialized", "d")]},
    ]
    add_code(s, MX, 3.04, 5.95, before)
    # after panel
    rect(s, MX + 6.18, 2.58, 5.95, 0.46, GREEN_SOFT, radius=0.09)
    add_icon(s, MX + 6.36, 2.70, 0.20, "i-check-circle", "green")
    txt(s, MX + 6.64, 2.70, 5.0, 0.3, [_p("With CodeGuard", s=12, b=True, c=GREEN)])
    after = [
        {"runs": [("raw_data = request.get_data(as_text=", "d"), ("True", "k"), (")", "d")]},
        {"runs": [("if not", "k"), (" raw_data: abort(", "d"), ("400", "n"), (")", "d")], "kind": "safe"},
        {"runs": [("decoded = base64.b64decode(raw_data, validate=", "d"), ("True", "k"), (")", "d")]},
        {"runs": [("data = json.loads(decoded.decode(", "d"), ("\"utf-8\"", "s"), ("))", "d")], "kind": "safe"},
        {"runs": [("# ← safe JSON instead of pickle", "c")], "kind": "safe"},
        {"runs": [("return ", "k"), ("jsonify(data)", "d")]},
    ]
    add_code(s, MX + 6.18, 3.04, 5.95, after)
    txt(s, MX, 4.65, 12.13, 0.6,
        [_p("The baseline takes the shortest path (", s=12, c=MUTED),
         _p("pickle", s=12, b=True, c=INK2, font=MONO),
         _p("). CodeGuard swaps in safe JSON, adds input validation and ", s=12, c=MUTED),
         _p("validate=True", s=12, font=MONO, c=INK2),
         _p(".", s=12, c=MUTED)], align=PP_ALIGN.CENTER)

    # ----------------------------------------------------- 12 · TEAM ADOPTION
    s = prs.slides.add_slide(blank)
    chrome(s, "Adoption · 11", 12)
    header(s, "Adoption", "Team Adoption",
           "One repo. One setup script. Works across every major AI coding agent.")
    agents = [("i-terminal", "Pi", "./scripts/setup-codeguard.sh"),
              ("i-sparkles", "Claude Code", "/plugin install codeguard-security@project-codeguard"),
              ("i-cpu", "Cursor", "codeguard-cursor.zip → copy to project"),
              ("i-git", "GitHub Copilot", "codeguard-copilot.zip → copy .github/"),
              ("i-globe", "Windsurf", "codeguard-windsurf.zip → copy to project"),
              ("i-box", "OpenAI Codex", "$skill-installer install …/project-codeguard")]
    # header row
    rect(s, MX, 2.60, 12.13, 0.42, SURFACE2, radius=0.06)
    txt(s, MX + 0.2, 2.68, 3.5, 0.25, [_p("AGENT", s=9.5, b=True, c=MUTED, spc=100)])
    txt(s, MX + 3.9, 2.68, 7.0, 0.25, [_p("INSTALL", s=9.5, b=True, c=MUTED, spc=100)])
    txt(s, MX + 11.0, 2.68, 1.0, 0.25, [_p("SETUP", s=9.5, b=True, c=MUTED, spc=100)])
    for i, (ic, agent, how) in enumerate(agents):
        t = 3.08 + i * 0.55
        if i % 2 == 0:
            rect(s, MX, t, 12.13, 0.52, WHITE, border=LINE)
        add_icon(s, MX + 0.2, t + 0.13, 0.24, ic, "red")
        txt(s, MX + 0.55, t + 0.14, 3.2, 0.25, [_p(agent, s=12, b=True, c=INK)])
        txt(s, MX + 3.9, t + 0.16, 6.9, 0.25, [_p(how, s=10.5, c=INK2, font=MONO)])
        rect(s, MX + 11.0, t + 0.12, 0.95, 0.28, GREEN_SOFT, radius=0.14)
        txt(s, MX + 11.0, t + 0.16, 0.95, 0.2, [_p("10 sec" if i == 0 else ("15 sec" if i == 1 else "30 sec"), s=9.5, b=True, c=GREEN, align=PP_ALIGN.CENTER)])
    txt(s, MX, 6.45, 12.13, 0.3,
        [_p("Skills, rules, AGENTS.md and the installer live in version control — team members clone and run one command.", s=10.5, c=FAINT)],
        align=PP_ALIGN.CENTER)

    # ----------------------------------------------------- 13 · TAKEAWAYS
    s = prs.slides.add_slide(blank)
    chrome(s, "Summary · 12", 13)
    header(s, "Summary", "Key Takeaways", None)
    takes = [("red", "Speed without structure is dangerous", "AI agents write code 2–5× faster — including vulnerabilities at the same rate. Without structure, you ship bugs faster."),
             ("cyan", "The agentic workflow is the foundation", "grill → spec → tickets → implement → review. Every output is verified; the AI is a teammate, not an oracle."),
             ("amber", "Security must be a separate axis", "Code review checks standards and spec. Security needs its own review with domain rules — don't let XSS compete with \u201crename this variable.\u201d"),
             ("green", "Prevention + detection = defense in depth", "CodeGuard during generation prevents; after generation it catches what slips through. 36.4% fewer vulnerabilities is the proof.")]
    for i, (cc, ti, de) in enumerate(takes):
        col, row = i % 2, i // 2
        l = MX + col * 6.18
        t = 2.55 + row * 1.95
        rect(s, l, t, 5.95, 1.75, SURFACE, radius=0.09, border=LINE)
        rect(s, l, t, 0.05, 1.75, COLOR_RGB[cc])
        txt(s, l + 0.22, t + 0.16, 5.5, 0.3, [_p(ti, s=14, b=True, c=INK)])
        txt(s, l + 0.22, t + 0.50, 5.55, 1.1, [_p(de, s=11.5, c=INK2, line_spacing=1.15)])
    txt(s, MX, 6.45, 12.13, 0.4,
        [_p("Secure by default. ", s=20, b=True, c=RED),
         _p("Not \u201csecure if we remember.\u201d", s=20, c=MUTED)], align=PP_ALIGN.CENTER)

    # ----------------------------------------------------- 14 · RESOURCES
    s = prs.slides.add_slide(blank)
    chrome(s, "Close · 13", 14)
    header(s, "Close", "Resources & Next Steps", None)
    res = [("i-db", "red", "This project", "github.com/abiolaks/security_skills_ai_coding", "Skills, rules, setup script and docs — everything from this talk."),
           ("i-shield", "cyan", "Project CodeGuard", "project-codeguard.org", "CoSAI / OASIS open standard. 23 rules, all major agents supported."),
           ("i-book", "green", "The study", "community.cisco.com", "\u201cCan Security Rules Make AI Generated Code Safer?\u201d — 2,717 prompts, GPT-5, 36.4% reduction."),
           ("i-refresh", "amber", "Agentic workflow", "github.com/cosai-oasis/project-codeguard", "Full source, MCP server and rule-translator tooling.")]
    for i, (ic, cc, ti, link, de) in enumerate(res):
        col, row = i % 2, i // 2
        l = MX + col * 6.18
        t = 2.55 + row * 1.75
        card(s, l, t, 5.95, 1.6, icon=ic, icon_color=cc, title=ti,
             lines=[_p(link, s=11, b=True, c=CYAN, font=MONO, sp=3), _p(de, s=10.5, c=INK2)])
    txt(s, MX, 6.20, 12.13, 0.6, [_p("Thank you. ", s=26, b=True, c=INK),
                                   _p("Questions?", s=26, c=MUTED)], align=PP_ALIGN.CENTER)

    prs.save(out_path)
    print(f"Saved deck to: {out_path}")
    print(f"Slides: {len(prs.slides._sldIdLst)}")


if __name__ == "__main__":
    out = "/Users/abiolaks/workspace/security_skills_ai_coding/slides/Secure_by_Default_AI_Agentic_Workflows.pptx"
    build(out)
